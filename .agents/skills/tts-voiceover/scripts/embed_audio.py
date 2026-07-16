#!/usr/bin/env python3
# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Embed per-slide WAV voice-over files into a PowerPoint deck.

Reads slide-NNN.wav files from an audio directory and adds them as embedded
media objects in the corresponding slides of a PPTX file. Adds animation
timing XML so PowerPoint recognizes the audio as narrations, enabling
'Use Recorded Timings and Narrations' in File > Export > Create a Video.

Usage:
    python embed_audio.py --input deck.pptx --audio-dir voice-over
    python embed_audio.py --input deck.pptx --audio-dir voice-over \
        --output deck-narrated.pptx
"""

from __future__ import annotations

import argparse
import logging
import sys
import wave
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Inches

logger = logging.getLogger(__name__)

EXIT_SUCCESS = 0
EXIT_FAILURE = 1
EXIT_ERROR = 2

AUDIO_MIME_TYPE = "audio/wav"
ICON_SIZE = Inches(0.1)
TIMING_BUFFER_MS = 1500

_PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_TIMING_TEMPLATE = (
    f'<p:timing xmlns:p="{_PPTX_NS}">'
    "<p:tnLst><p:par>"
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    "<p:childTnLst>"
    '<p:seq concurrent="1" nextAc="seek">'
    '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
    "<p:childTnLst><p:par>"
    '<p:cTn id="3" fill="hold">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    "<p:childTnLst><p:par>"
    '<p:cTn id="4" fill="hold">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    "<p:childTnLst>"
    '<p:cmd type="call" cmd="playFrom(0)"><p:cBhvr>'
    '<p:cTn id="5" dur="0" fill="hold"/>'
    '<p:tgtEl><p:spTgt spid="0"/></p:tgtEl>'
    "</p:cBhvr></p:cmd>"
    "</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>"
    "</p:childTnLst></p:cTn>"
    "<p:prevCondLst>"
    '<p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
    "</p:prevCondLst>"
    "<p:nextCondLst>"
    '<p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
    "</p:nextCondLst>"
    "</p:seq></p:childTnLst></p:cTn>"
    "</p:par></p:tnLst></p:timing>"
)


def get_wav_duration_ms(wav_path: Path) -> int:
    """Return WAV file duration in milliseconds with buffer.

    Args:
        wav_path: Path to the WAV audio file.

    Returns:
        Duration in milliseconds plus ``TIMING_BUFFER_MS``.
    """
    with wave.open(str(wav_path), "rb") as wf:
        frames = wf.getnframes()
        rate = wf.getframerate()
        return int((frames / float(rate)) * 1000) + TIMING_BUFFER_MS


def _add_narration_timing(slide: Slide, shape_id: int, duration_ms: int) -> None:
    """Add auto-play narration timing XML to a slide.

    Creates the ``p:timing`` element structure that PowerPoint generates
    when using Record Slide Show, enabling 'Use Recorded Timings and
    Narrations' in video export.

    Args:
        slide: The slide to modify.
        shape_id: The ``spId`` of the embedded audio shape.
        duration_ms: Audio duration in milliseconds.
    """
    existing = slide._element.find(qn("p:timing"))
    if existing is not None:
        # Warn whenever an existing timing element is replaced, since any
        # authored animation (entrance effect, click sequence) produces at
        # least one p:seq that will be overwritten.
        child_seqs = existing.findall(f".//{qn('p:seq')}")
        if child_seqs:
            logger.warning(
                "Replacing existing slide timing (%d sequence(s)) for shape %d; "
                "authored animations on this slide will be overwritten.",
                len(child_seqs),
                shape_id,
            )
        slide._element.remove(existing)

    parser = etree.XMLParser(resolve_entities=False, no_network=True)
    timing = etree.fromstring(_TIMING_TEMPLATE, parser)
    ns = {"p": _PPTX_NS}
    sp_tgt = timing.find(".//p:spTgt", ns)
    if sp_tgt is not None:
        sp_tgt.set("spid", str(shape_id))
    else:
        logger.warning(
            "spTgt element not found in timing template for shape %d; "
            "audio shape link will be missing.",
            shape_id,
        )
    ctn_dur = timing.find(".//p:cTn[@id='5']", ns)
    if ctn_dur is not None:
        ctn_dur.set("dur", str(duration_ms))
    else:
        logger.warning(
            "cTn[@id='5'] not found in timing template for shape %d; "
            "audio duration will be unset.",
            shape_id,
        )
    slide._element.append(timing)


def _set_slide_transition(slide: Slide, duration_ms: int) -> None:
    """Set slide auto-advance timing after audio duration.

    Sets ``advClick="0"`` so slides advance only on the audio timer,
    not on manual click. To re-enable click-to-advance after embedding,
    use the Transitions tab in PowerPoint.

    Args:
        slide: The slide to modify.
        duration_ms: Auto-advance delay in milliseconds.
    """
    existing = slide._element.find(qn("p:transition"))
    if existing is not None:
        slide._element.remove(existing)

    # advClick="0" prevents accidental click-to-skip during audio playback;
    # slides advance only when the audio timer expires.
    transition = slide._element.makeelement(
        qn("p:transition"),
        {"advClick": "0", "advTm": str(duration_ms)},
    )
    timing = slide._element.find(qn("p:timing"))
    if timing is not None:
        timing.addprevious(transition)
    else:
        slide._element.append(transition)


def embed_slide_audio(slide: Slide, wav_path: Path) -> bool:
    """Embed a WAV file into a slide as a media object.

    Adds narration timing XML and slide auto-advance so PowerPoint
    recognizes the audio for video export.

    Args:
        slide: The target slide.
        wav_path: Path to the WAV audio file to embed.

    Returns:
        ``True`` on success, ``False`` on failure.
    """
    try:
        movie_shape = slide.shapes.add_movie(
            str(wav_path),
            left=0,
            top=0,
            width=ICON_SIZE,
            height=ICON_SIZE,
            mime_type=AUDIO_MIME_TYPE,
        )
        shape_id: int = movie_shape.shape_id
        duration_ms = get_wav_duration_ms(wav_path)
        _add_narration_timing(slide, shape_id, duration_ms)
        _set_slide_transition(slide, duration_ms)
        return True
    except Exception as exc:  # python-pptx raises varied internal exceptions
        logger.exception(
            "Failed to embed audio %s (%s)", wav_path.name, type(exc).__name__
        )
        return False


def create_parser() -> argparse.ArgumentParser:
    """Create and configure the argument parser."""
    parser = argparse.ArgumentParser(
        description="Embed per-slide WAV voice-over files into a PPTX deck"
    )
    parser.add_argument(
        "--input",
        type=Path,
        required=True,
        help="Source PPTX file path",
    )
    parser.add_argument(
        "--audio-dir",
        type=Path,
        default=Path("voice-over"),
        help="Directory containing slide-NNN.wav files (default: voice-over)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Output PPTX file path (default: input stem + '-narrated.pptx')",
    )
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable verbose (DEBUG) logging",
    )
    return parser


def configure_logging(verbose: bool = False) -> None:
    """Configure logging based on verbosity level."""
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=level, format="%(levelname)s: %(message)s")


def _run(args: argparse.Namespace) -> int:
    """Execute audio embedding logic."""

    input_path: Path = args.input
    audio_dir: Path = args.audio_dir

    if not input_path.is_file():
        logger.error("Input PPTX not found: %s", input_path)
        return EXIT_FAILURE

    if not audio_dir.is_dir():
        logger.error("Audio directory not found: %s", audio_dir)
        return EXIT_FAILURE

    output_path: Path = args.output or input_path.with_name(
        f"{input_path.stem}-narrated.pptx"
    )

    if output_path.resolve() == input_path.resolve():
        logger.error(
            "Output path must differ from input path to avoid overwriting the source"
        )
        return EXIT_ERROR

    prs = Presentation(str(input_path))
    embedded_count = 0
    failed_count = 0

    # Build a mapping from slide number to WAV path so embedding matches
    # the directory names used by generate_voiceover.py rather than
    # re-deriving names from the enumerate index.
    wav_files: dict[int, Path] = {}
    for wav in sorted(audio_dir.glob("slide-*.wav")):
        try:
            num = int(wav.stem.split("-")[1])
            wav_files[num] = wav
        except (IndexError, ValueError):
            logger.warning("Ignoring unexpected file: %s", wav.name)

    for idx, slide in enumerate(prs.slides, start=1):
        wav_path = wav_files.get(idx)
        if wav_path is None:
            logger.info("SKIP slide %d: no WAV found", idx)
            continue

        if embed_slide_audio(slide, wav_path):
            embedded_count += 1
            logger.info("Embedded %s into slide %d", wav_path.name, idx)
        else:
            logger.error("FAILED to embed %s into slide %d", wav_path.name, idx)
            failed_count += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if embedded_count == 0:
        logger.error(
            "No audio files were embedded. Verify that slide-NNN.wav files exist in %s",
            audio_dir,
        )
        return EXIT_FAILURE

    try:
        prs.save(str(output_path))
    except OSError as exc:
        logger.error("Failed to save output PPTX %s: %s", output_path, exc)
        return EXIT_FAILURE

    logger.info("Saved %s with %d embedded audio files", output_path, embedded_count)

    if failed_count > 0:
        logger.warning(
            "Completed with %d failure(s); %d slide(s) embedded successfully.",
            failed_count,
            embedded_count,
        )
        return EXIT_FAILURE
    return EXIT_SUCCESS


def main() -> int:
    """Entry point for audio embedding."""
    parser = create_parser()
    args = parser.parse_args()
    configure_logging(verbose=args.verbose)
    try:
        return _run(args)
    except KeyboardInterrupt:
        return 130
    except BrokenPipeError:
        sys.stderr.close()
        return 1


if __name__ == "__main__":
    sys.exit(main())
