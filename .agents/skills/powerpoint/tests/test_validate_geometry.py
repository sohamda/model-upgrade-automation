# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Tests for validate_geometry module."""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from validate_geometry import (
    _is_accent_bar,
    _shape_label,
    check_adjacent_gaps,
    check_boundary_overflow,
    check_edge_margins,
    check_title_clearance,
    create_parser,
    generate_report,
    main,
    max_severity,
    validate_geometry,
    validate_slide_geometry,
)


@pytest.fixture()
def simple_deck(tmp_path):
    """Create a minimal PPTX with 2 slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def deck_with_shapes(tmp_path):
    """PPTX with shapes at known positions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Accent bar at top (should be exempted)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.05),
    )
    bar.name = "accent_bar"

    # Title at proper position
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.0), Inches(0.7)
    )
    title.name = "title"
    title.text_frame.text = "Slide Title"

    # Content below title
    content = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(11.0), Inches(4.0)
    )
    content.name = "content"
    content.text_frame.text = "Content area"

    path = tmp_path / "shapes.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def deck_with_violations(tmp_path):
    """PPTX with deliberate margin and overflow violations."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Shape too close to left edge (0.2" < 0.5")
    tight = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.2),
        Inches(0.5),
        Inches(2.0),
        Inches(1.0),
    )
    tight.name = "tight_left"

    # Shape overflowing right boundary
    overflow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(12.0),
        Inches(1.0),
        Inches(2.0),
        Inches(1.0),
    )
    overflow.name = "overflow_right"

    path = tmp_path / "violations.pptx"
    prs.save(str(path))
    return path


class TestIsAccentBar:
    """Tests for _is_accent_bar."""

    def test_full_width_thin_bar(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.05),
        )
        assert _is_accent_bar(shape, 13.333) is True

    def test_regular_shape_not_bar(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        assert _is_accent_bar(shape, 13.333) is False

    def test_tall_shape_not_bar(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.5),
        )
        assert _is_accent_bar(shape, 13.333) is False


class TestShapeLabel:
    """Tests for _shape_label."""

    def test_named_shape_with_text(self, blank_slide):
        tb = blank_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.name = "Title 1"
        tb.text_frame.text = "Hello World"
        label = _shape_label(tb)
        assert "Title 1" in label
        assert "Hello World" in label

    def test_named_shape_without_text(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        shape.name = "rect1"
        assert _shape_label(shape) == "rect1"


class TestCheckBoundaryOverflow:
    """Tests for check_boundary_overflow."""

    def test_no_overflow(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        issues = check_boundary_overflow(shape, 13.333, 7.5)
        assert len(issues) == 0

    def test_right_overflow(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(12),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        issues = check_boundary_overflow(shape, 13.333, 7.5)
        assert len(issues) == 1
        assert issues[0]["check_type"] == "boundary_overflow"
        assert issues[0]["severity"] == "error"

    def test_bottom_overflow(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(7),
            Inches(2),
            Inches(1),
        )
        issues = check_boundary_overflow(shape, 13.333, 7.5)
        assert len(issues) == 1
        assert "bottom" in issues[0]["description"].lower()

    def test_left_overflow(self, blank_slide):
        """Shape with negative left position is detected."""
        from pptx.util import Emu

        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(-91440),  # -0.1 inches
            Inches(1),
            Inches(2),
            Inches(1),
        )
        issues = check_boundary_overflow(shape, 13.333, 7.5)
        assert any("left" in i["description"].lower() for i in issues)

    def test_top_overflow(self, blank_slide):
        """Shape with negative top position is detected."""
        from pptx.util import Emu

        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Emu(-91440),  # -0.1 inches
            Inches(2),
            Inches(1),
        )
        issues = check_boundary_overflow(shape, 13.333, 7.5)
        assert any("top" in i["description"].lower() for i in issues)


class TestCheckEdgeMargins:
    """Tests for check_edge_margins."""

    def test_within_margins(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8),
            Inches(0.8),
            Inches(4),
            Inches(2),
        )
        issues = check_edge_margins(shape, 13.333, 7.5, 0.5)
        assert len(issues) == 0

    def test_too_close_to_left(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        issues = check_edge_margins(shape, 13.333, 7.5, 0.5)
        assert any(i["check_type"] == "edge_margin" for i in issues)

    def test_too_close_to_top(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(0.2),
            Inches(2),
            Inches(1),
        )
        issues = check_edge_margins(shape, 13.333, 7.5, 0.5)
        assert any("top" in i["description"].lower() for i in issues)


class TestCheckAdjacentGaps:
    """Tests for check_adjacent_gaps."""

    def test_sufficient_gap(self, blank_slide):
        s1 = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        s2 = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(2.5),
            Inches(4),
            Inches(1),
        )
        issues = check_adjacent_gaps([s1, s2], 0.3)
        assert len(issues) == 0

    def test_insufficient_gap(self, blank_slide):
        s1 = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        s2 = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(2.1),
            Inches(4),
            Inches(1),
        )
        issues = check_adjacent_gaps([s1, s2], 0.3)
        assert len(issues) == 1
        assert issues[0]["check_type"] == "adjacent_gap"


class TestCheckTitleClearance:
    """Tests for check_title_clearance."""

    def test_sufficient_clearance(self, blank_slide):
        title = blank_slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(10), Inches(0.7)
        )
        title.name = "title"
        content = blank_slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(10), Inches(4)
        )
        content.name = "content"
        issues = check_title_clearance([title, content], 0.2)
        assert len(issues) == 0

    def test_tight_clearance(self, blank_slide):
        title = blank_slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(10), Inches(0.7)
        )
        title.name = "title"
        content = blank_slide.shapes.add_textbox(
            Inches(1), Inches(1.25), Inches(10), Inches(4)
        )
        content.name = "content"
        issues = check_title_clearance([title, content], 0.2)
        assert len(issues) == 1
        assert issues[0]["check_type"] == "title_clearance"


class TestValidateSlideGeometry:
    """Tests for validate_slide_geometry."""

    def test_clean_slide(self, blank_slide):
        result = validate_slide_geometry(
            blank_slide, 1, 13.333, 7.5, margin=0.5, gap=0.3, clearance=0.2
        )
        assert result["slide_number"] == 1
        assert result["overall_quality"] == "good"

    def test_slide_with_issues(self, blank_slide):
        blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.1),
            Inches(0.1),
            Inches(2),
            Inches(1),
        )
        result = validate_slide_geometry(
            blank_slide, 1, 13.333, 7.5, margin=0.5, gap=0.3, clearance=0.2
        )
        assert result["overall_quality"] == "needs-attention"
        assert len(result["issues"]) > 0


class TestValidateGeometry:
    """Tests for validate_geometry."""

    def test_full_validation(self, simple_deck):
        results = validate_geometry(simple_deck)
        assert results["source"] == "geometry-validation"
        assert results["slide_count"] == 2
        assert len(results["slides"]) == 2

    def test_slide_filter(self, simple_deck):
        results = validate_geometry(simple_deck, slide_filter={1})
        assert len(results["slides"]) == 1
        assert results["slides"][0]["slide_number"] == 1

    def test_clean_deck(self, deck_with_shapes):
        results = validate_geometry(deck_with_shapes)
        # Deck has shapes well within boundaries; may have minor
        # warnings from right-edge proximity of 11" wide content
        sev = max_severity(results)
        assert sev in ("none", "info", "warning")


class TestGenerateReport:
    """Tests for generate_report."""

    def test_report_structure(self, simple_deck):
        results = validate_geometry(simple_deck)
        report = generate_report(results)
        assert "# Geometry Validation Report" in report
        assert "## Summary" in report
        assert "## Per-Slide Findings" in report

    def test_report_with_issues(self, deck_with_violations):
        results = validate_geometry(deck_with_violations)
        report = generate_report(results)
        assert "warning" in report.lower() or "error" in report.lower()


class TestMaxSeverity:
    """Tests for max_severity."""

    def test_no_issues(self):
        results = {"slides": [{"issues": []}]}
        assert max_severity(results) == "none"

    def test_error_dominates(self):
        results = {
            "slides": [
                {
                    "issues": [
                        {"severity": "info"},
                        {"severity": "error"},
                        {"severity": "warning"},
                    ]
                }
            ]
        }
        assert max_severity(results) == "error"

    def test_warning_over_info(self):
        results = {
            "slides": [{"issues": [{"severity": "info"}, {"severity": "warning"}]}]
        }
        assert max_severity(results) == "warning"


class TestCreateParser:
    """Tests for create_parser."""

    def test_required_input(self):
        parser = create_parser()
        args = parser.parse_args(["--input", "test.pptx"])
        assert str(args.input) == "test.pptx"

    def test_defaults(self):
        parser = create_parser()
        args = parser.parse_args(["--input", "test.pptx"])
        assert args.margin == 0.5
        assert args.gap == 0.3
        assert args.clearance == 0.2

    def test_custom_thresholds(self):
        parser = create_parser()
        args = parser.parse_args(
            ["--input", "t.pptx", "--margin", "0.6", "--gap", "0.4"]
        )
        assert args.margin == 0.6
        assert args.gap == 0.4


class TestMain:
    """Tests for main entry point."""

    def test_valid_deck(self, simple_deck, monkeypatch):
        monkeypatch.setattr(
            "sys.argv",
            ["validate_geometry", "--input", str(simple_deck)],
        )
        rc = main()
        assert rc == 0

    def test_missing_file(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "sys.argv",
            ["validate_geometry", "--input", str(tmp_path / "missing.pptx")],
        )
        rc = main()
        assert rc == 2

    def test_json_output(self, simple_deck, tmp_path, monkeypatch):
        out = tmp_path / "results.json"
        monkeypatch.setattr(
            "sys.argv",
            [
                "validate_geometry",
                "--input",
                str(simple_deck),
                "--output",
                str(out),
            ],
        )
        main()
        assert out.exists()
        data = json.loads(out.read_text())
        assert "slides" in data

    def test_report_output(self, simple_deck, tmp_path, monkeypatch):
        report = tmp_path / "report.md"
        monkeypatch.setattr(
            "sys.argv",
            [
                "validate_geometry",
                "--input",
                str(simple_deck),
                "--report",
                str(report),
            ],
        )
        main()
        assert report.exists()
        assert "# Geometry Validation Report" in report.read_text(encoding="utf-8")

    def test_per_slide_dir(self, simple_deck, tmp_path, monkeypatch):
        per_slide = tmp_path / "per-slide"
        monkeypatch.setattr(
            "sys.argv",
            [
                "validate_geometry",
                "--input",
                str(simple_deck),
                "--per-slide-dir",
                str(per_slide),
            ],
        )
        main()
        assert per_slide.exists()
        geom_files = list(per_slide.glob("slide-*-geometry.json"))
        assert len(geom_files) == 2
