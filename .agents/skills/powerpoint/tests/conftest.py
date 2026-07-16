# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Shared fixtures for PowerPoint skill tests."""

import io
import os
import struct
import zlib
from pathlib import Path

import pytest
from hypothesis import HealthCheck, settings
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Hypothesis profiles
settings.register_profile(
    "ci",
    max_examples=200,
    derandomize=True,
    deadline=None,
    database=None,
    print_blob=True,
    suppress_health_check=[HealthCheck.too_slow],
)
settings.register_profile(
    "dev",
    max_examples=50,
    deadline=None,
)
settings.load_profile("ci" if os.environ.get("CI") else "dev")


# Plain functions callable from @given-decorated Hypothesis tests,
# which cannot use pytest fixtures.


def make_blank_presentation():
    """Create a fresh Presentation with standard widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def make_blank_slide():
    """Create a blank slide on a fresh presentation."""
    prs = make_blank_presentation()
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _minimal_png_bytes() -> bytes:
    """Create a minimal valid 1x1 red PNG in memory."""

    def _chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw_row = b"\x00\xff\x00\x00"  # filter byte + RGB
    idat_data = zlib.compress(raw_row)
    return (
        signature
        + _chunk(b"IHDR", ihdr_data)
        + _chunk(b"IDAT", idat_data)
        + _chunk(b"IEND", b"")
    )


def _apply_srgb_color(
    clr_scheme: etree._Element,
    ns: dict[str, str],
    color_name: str,
    hex_val: str,
) -> None:
    """Set a theme color's srgbClr value, raising if the node is absent."""
    node = clr_scheme.find(f"a:{color_name}", ns)
    if node is None:
        raise ValueError(
            f"Could not find theme color node a:{color_name} in clrScheme."
        )
    for child in list(node):
        node.remove(child)
    srgb = etree.SubElement(
        node,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
    )
    srgb.set("val", hex_val)


def _set_theme_colors(prs: Presentation) -> None:
    """Set specific theme colors via the theme part's public
    blob setter."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    slide_master_part = prs.slide_masters[0].part
    theme_part = None
    for rel in slide_master_part.rels.values():
        if "theme" in rel.reltype:
            theme_part = rel.target_part
            break

    if theme_part is None:
        raise ValueError("Could not find theme part in slide master relationships.")

    theme_element = etree.fromstring(theme_part.blob)

    clr_scheme = theme_element.find(".//a:clrScheme", ns)
    if clr_scheme is None:
        raise ValueError("Could not find clrScheme in theme XML.")

    _apply_srgb_color(clr_scheme, ns, "dk1", "000000")
    _apply_srgb_color(clr_scheme, ns, "accent1", "4F81BD")

    theme_part.blob = etree.tostring(
        theme_element,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )


def generate_minimal_fixture(output_path: Path) -> None:
    """Builds the minimal test PPTX programmatically."""
    prs = make_blank_presentation()

    prs.core_properties.title = "Minimal Test Fixture"
    prs.core_properties.author = "HVE Core Test Fixture"

    _set_theme_colors(prs)

    slide_layout_1 = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout_1)
    slide1.placeholders[0].text = "Test Fixture Presentation"
    slide1.placeholders[1].text = "Slide with theme colors and notes"

    title_shape = slide1.placeholders[0]
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    notes_slide = slide1.notes_slide
    notes_slide.notes_text_frame.text = "This is a speaker note for slide 1."

    slide_layout_2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout_2)
    slide2.placeholders[0].text = "Slide with Image"
    slide2.placeholders[1].text = "Below is an embedded image."
    slide2.notes_slide.notes_text_frame.text = "This is a speaker note for slide 2."
    image_stream = io.BytesIO(_minimal_png_bytes())
    slide2.shapes.add_picture(image_stream, Inches(1), Inches(2), width=Inches(2))

    prs.save(str(output_path))


# Fixtures


@pytest.fixture()
def blank_presentation():
    """Fresh Presentation with standard widescreen dimensions."""
    return make_blank_presentation()


@pytest.fixture()
def blank_slide(blank_presentation):
    """Blank slide added to a fresh presentation."""
    layout = blank_presentation.slide_layouts[6]  # Blank layout
    return blank_presentation.slides.add_slide(layout)


@pytest.fixture()
def sample_textbox(blank_slide):
    """Slide with a textbox containing known text and formatting."""
    txBox = blank_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.text = "Sample Text"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True
    return txBox


@pytest.fixture()
def sample_shape(blank_slide):
    """Slide with a rectangle shape having fill and text."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = blank_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2),
        Inches(2),
        Inches(3),
        Inches(2),
    )
    shape.text = "Shape Text"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
    return shape


@pytest.fixture()
def sample_image_path(tmp_path):
    """Minimal valid PNG file at a temporary path."""
    img = tmp_path / "test.png"
    img.write_bytes(_minimal_png_bytes())
    return img


@pytest.fixture(scope="session")
def powerpoint_fixture_dir() -> Path:
    return Path(__file__).parent / "fixtures"


@pytest.fixture(scope="session")
def minimal_test_fixture_path(
    tmp_path_factory: pytest.TempPathFactory,
) -> Path:
    """Generates the minimal test fixture on the fly and returns its path."""
    fixture_dir = tmp_path_factory.mktemp("fixtures")
    pptx_path = fixture_dir / "minimal_test_fixture.pptx"
    generate_minimal_fixture(pptx_path)
    return pptx_path


@pytest.fixture(scope="session")
def malformed_pdf_dir(powerpoint_fixture_dir: Path) -> Path:
    """Directory holding tiny on-disk malformed-PDF fixtures."""
    return powerpoint_fixture_dir / "malformed"


@pytest.fixture()
def minimal_valid_pdf(tmp_path):
    """Factory writing bytes that pass ``validate_pdf_path`` magic+size checks.

    The bytes after the magic prefix are not a parsable PDF; tests using
    this fixture must mock ``fitz.open`` to bypass real C-level parsing.
    """

    def _make(name: str = "minimal.pdf") -> Path:
        pdf = tmp_path / name
        pdf.write_bytes(b"%PDF-1.4\n%fake\n")
        return pdf

    return _make


@pytest.fixture()
def oversized_pdf(tmp_path):
    """Factory writing a >MAX_PDF_BYTES file via sparse file (no large commit).

    Returns a Path whose ``stat().st_size`` exceeds the configured ceiling
    while occupying near-zero bytes on disk. The caller may pass a custom
    ``max_bytes`` to ``validate_pdf_path`` to keep the synthetic size
    cheap; defaults target the production ceiling.
    """

    def _make(size_bytes: int = 100 * 1024 * 1024 + 1, name: str = "huge.pdf") -> Path:
        pdf = tmp_path / name
        with pdf.open("wb") as fh:
            fh.write(b"%PDF-1.4\n")
            # Sparse file: seek beyond and write a single byte. This makes
            # st_size large without committing megabytes of zeros to disk
            # on filesystems that support sparse files (APFS, ext4, NTFS).
            fh.seek(size_bytes - 1)
            fh.write(b"\x00")
        return pdf

    return _make


@pytest.fixture()
def many_page_pdf(tmp_path):
    """Factory writing a magic-valid file paired with a fitz mock for N pages.

    Returns a tuple ``(path, page_count)``. The on-disk bytes are not a real
    PDF; the caller is expected to mock ``fitz.open`` so that ``len(doc)``
    returns ``page_count`` to drive the page-count ceiling check in
    :func:`pdf_safety.safe_open_pdf`.
    """

    def _make(page_count: int = 1001, name: str = "many.pdf") -> tuple[Path, int]:
        pdf = tmp_path / name
        pdf.write_bytes(b"%PDF-1.4\n%fake\n")
        return pdf, page_count

    return _make
