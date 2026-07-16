# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Regression tests for fixture generation behavior."""

from pathlib import Path

import pytest
from conftest import generate_minimal_fixture
from lxml import etree
from pptx import Presentation


def test_apply_srgb_color_raises_on_missing_node() -> None:
    """Missing theme color nodes should raise a clear error."""
    clr_scheme = etree.Element(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme"
    )
    lt1 = etree.SubElement(
        clr_scheme,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}lt1",
    )
    etree.SubElement(
        lt1, "{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr"
    )

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    with pytest.raises(ValueError, match="Could not find theme color node a:dk1"):
        from conftest import _apply_srgb_color

        _apply_srgb_color(clr_scheme, ns, "dk1", "000000")


@pytest.mark.integration
def test_generated_fixture_has_notes_on_all_slides(tmp_path: Path) -> None:
    """Every slide in the generated fixture should include non-empty notes."""
    pptx_path = tmp_path / "fixture.pptx"

    generate_minimal_fixture(pptx_path)

    prs = Presentation(str(pptx_path))

    for slide in prs.slides:
        assert slide.has_notes_slide
        assert slide.notes_slide.notes_text_frame.text.strip()
