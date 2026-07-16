# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Tests for embed_audio module."""

from __future__ import annotations

import struct

import pytest
from embed_audio import (
    AUDIO_PATTERN,
    create_parser,
    create_poster_frame,
    discover_audio_files,
    embed_audio,
    main,
    run,
)
from pptx import Presentation
from pptx.util import Inches


def _make_wav_bytes(duration_ms: int = 100) -> bytes:
    """Create minimal valid WAV file bytes."""
    sample_rate = 16000
    num_samples = int(sample_rate * duration_ms / 1000)
    data = b"\x00\x00" * num_samples
    data_size = len(data)
    fmt_size = 16
    file_size = 4 + (8 + fmt_size) + (8 + data_size)
    header = struct.pack(
        "<4sI4s4sIHHIIHH4sI",
        b"RIFF",
        file_size,
        b"WAVE",
        b"fmt ",
        fmt_size,
        1,  # PCM
        1,  # mono
        sample_rate,
        sample_rate * 2,  # byte rate
        2,  # block align
        16,  # bits per sample
        b"data",
        data_size,
    )
    return header + data


@pytest.fixture()
def simple_deck(tmp_path):
    """Create a minimal 3-slide PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    for _ in range(3):
        prs.slides.add_slide(layout)
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def audio_dir(tmp_path):
    """Create a directory with 3 WAV files."""
    d = tmp_path / "audio"
    d.mkdir()
    wav = _make_wav_bytes()
    for i in range(1, 4):
        (d / f"slide-{i:03d}.wav").write_bytes(wav)
    return d


class TestAudioPattern:
    """Tests for the AUDIO_PATTERN regex."""

    def test_matches_standard(self):
        m = AUDIO_PATTERN.match("slide-001.wav")
        assert m is not None
        assert m.group(1) == "001"

    def test_matches_large_number(self):
        m = AUDIO_PATTERN.match("slide-123.wav")
        assert m is not None
        assert m.group(1) == "123"

    def test_rejects_wrong_prefix(self):
        assert AUDIO_PATTERN.match("audio-001.wav") is None

    def test_rejects_wrong_extension(self):
        assert AUDIO_PATTERN.match("slide-001.mp3") is None


class TestDiscoverAudioFiles:
    """Tests for discover_audio_files."""

    def test_finds_wav_files(self, audio_dir):
        mapping = discover_audio_files(audio_dir)
        assert len(mapping) == 3
        assert 1 in mapping
        assert 2 in mapping
        assert 3 in mapping

    def test_empty_dir(self, tmp_path):
        d = tmp_path / "empty"
        d.mkdir()
        assert discover_audio_files(d) == {}

    def test_ignores_non_wav(self, tmp_path):
        d = tmp_path / "mixed"
        d.mkdir()
        (d / "slide-001.wav").write_bytes(_make_wav_bytes())
        (d / "slide-002.mp3").write_bytes(b"\x00")
        (d / "README.md").write_text("hi")
        mapping = discover_audio_files(d)
        assert len(mapping) == 1
        assert 1 in mapping


class TestCreatePosterFrame:
    """Tests for create_poster_frame."""

    def test_creates_png_file(self):
        path = create_poster_frame()
        assert path.exists()
        assert path.suffix == ".png"
        data = path.read_bytes()
        assert data[:4] == b"\x89PNG"
        # Clean up temp file
        path.unlink(missing_ok=True)


class TestEmbedAudio:
    """Tests for embed_audio function."""

    def test_embeds_audio_on_slide(self, simple_deck, audio_dir, tmp_path):
        from pptx import Presentation as Prs

        output = tmp_path / "out.pptx"
        prs = Prs(str(simple_deck))
        poster = create_poster_frame()
        try:
            embedded = embed_audio(prs, {1: audio_dir / "slide-001.wav"}, None, poster)
        finally:
            poster.unlink(missing_ok=True)
        assert embedded == 1
        prs.save(str(output))
        assert output.exists()

    def test_slide_filter(self, simple_deck, audio_dir, tmp_path):
        from pptx import Presentation as Prs

        prs = Prs(str(simple_deck))
        audio_map = discover_audio_files(audio_dir)
        poster = create_poster_frame()
        try:
            embedded = embed_audio(prs, audio_map, {2}, poster)
        finally:
            poster.unlink(missing_ok=True)
        assert embedded == 1


class TestCreateParser:
    """Tests for create_parser."""

    def test_required_args(self):
        parser = create_parser()
        args = parser.parse_args(
            ["--input", "d.pptx", "--audio-dir", "a/", "--output", "o.pptx"]
        )
        assert str(args.input) == "d.pptx"

    def test_optional_slides(self):
        parser = create_parser()
        args = parser.parse_args(
            [
                "--input",
                "d.pptx",
                "--audio-dir",
                "a/",
                "--output",
                "o.pptx",
                "--slides",
                "1,3",
            ]
        )
        assert args.slides == "1,3"


class TestRun:
    """Tests for run function."""

    def test_full_embed(self, simple_deck, audio_dir, tmp_path):
        parser = create_parser()
        output = tmp_path / "narrated.pptx"
        args = parser.parse_args(
            [
                "--input",
                str(simple_deck),
                "--audio-dir",
                str(audio_dir),
                "--output",
                str(output),
            ]
        )
        rc = run(args)
        assert rc == 0
        assert output.exists()

    def test_missing_input(self, audio_dir, tmp_path):
        parser = create_parser()
        args = parser.parse_args(
            [
                "--input",
                str(tmp_path / "missing.pptx"),
                "--audio-dir",
                str(audio_dir),
                "--output",
                str(tmp_path / "out.pptx"),
            ]
        )
        rc = run(args)
        assert rc == 2

    def test_missing_audio_dir(self, simple_deck, tmp_path):
        parser = create_parser()
        args = parser.parse_args(
            [
                "--input",
                str(simple_deck),
                "--audio-dir",
                str(tmp_path / "no-audio"),
                "--output",
                str(tmp_path / "out.pptx"),
            ]
        )
        rc = run(args)
        assert rc == 2

    def test_no_matching_audio(self, simple_deck, tmp_path):
        empty_audio = tmp_path / "empty-audio"
        empty_audio.mkdir()
        parser = create_parser()
        args = parser.parse_args(
            [
                "--input",
                str(simple_deck),
                "--audio-dir",
                str(empty_audio),
                "--output",
                str(tmp_path / "out.pptx"),
            ]
        )
        rc = run(args)
        assert rc == 1  # EXIT_FAILURE: no matching audio files


class TestMain:
    """Tests for main entry point."""

    def test_success(self, simple_deck, audio_dir, tmp_path, monkeypatch):
        output = tmp_path / "main-out.pptx"
        monkeypatch.setattr(
            "sys.argv",
            [
                "embed_audio",
                "--input",
                str(simple_deck),
                "--audio-dir",
                str(audio_dir),
                "--output",
                str(output),
            ],
        )
        rc = main()
        assert rc == 0
        assert output.exists()
