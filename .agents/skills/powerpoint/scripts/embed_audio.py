#!/usr/bin/env python3
# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Embed WAV audio files into a PowerPoint deck, one per slide.

Matches audio files to slides by naming convention (slide-001.wav → slide 1)
and embeds each as an audio shape using python-pptx's add_movie API.

Usage::

    python embed_audio.py --input deck.pptx \
        --audio-dir voice-over/ --output out.pptx
    python embed_audio.py --input deck.pptx \
        --audio-dir voice-over/ --output out.pptx \
        --slides "1,3,5"
    python embed_audio.py --input deck.pptx \
        --audio-dir voice-over/ --output out.pptx -v
"""

from __future__ import annotations

import argparse
import io
import logging
import re
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pptx_utils import (
    EXIT_ERROR,
    EXIT_FAILURE,
    EXIT_SUCCESS,
    configure_logging,
    parse_slide_filter,
)

logger = logging.getLogger(__name__)

AUDIO_PATTERN = re.compile(r"^slide-(\d+)\.wav$", re.IGNORECASE)

AUDIO_LEFT = Inches(0.1)
AUDIO_WIDTH = Inches(0.3)
AUDIO_HEIGHT = Inches(0.3)
AUDIO_OFFSCREEN_OFFSET = Inches(0.5)


def create_parser() -> argparse.ArgumentParser:
    """Create and configure argument parser."""
    parser = argparse.ArgumentParser(
        description="Embed WAV audio files into a PowerPoint deck"
    )
    parser.add_argument(
        "--input", required=True, type=Path, help="Source PPTX file path"
    )
    parser.add_argument(
        "--audio-dir", required=True, type=Path, help="Directory containing WAV files"
    )
    parser.add_argument(
        "--output", required=True, type=Path, help="Output PPTX file path"
    )
    parser.add_argument(
        "--slides",
        help="Comma-separated slide numbers to embed audio on (1-based, default: all)",
    )
    parser.add_argument(
        "-v", "--verbose", action="store_true", help="Enable verbose output"
    )
    return parser


def discover_audio_files(audio_dir: Path) -> dict[int, Path]:
    """Map slide numbers to WAV file paths found in the audio directory.

    Scans for files matching the ``slide-NNN.wav`` naming convention.

    Args:
        audio_dir: Directory to scan for WAV files.

    Returns:
        Dictionary mapping 1-based slide numbers to their WAV file paths.
    """
    mapping: dict[int, Path] = {}
    for entry in sorted(audio_dir.iterdir()):
        if not entry.is_file():
            continue
        match = AUDIO_PATTERN.match(entry.name)
        if match:
            slide_num = int(match.group(1))
            mapping[slide_num] = entry
            logger.debug("Found audio for slide %d: %s", slide_num, entry.name)
    return mapping


def create_poster_frame() -> Path:
    """Create a minimal 1x1 transparent PNG for the audio poster frame.

    python-pptx's ``add_movie`` requires a poster frame image. This creates
    a temporary transparent PNG so the audio shape has no visible thumbnail.

    Returns:
        Path to the temporary PNG file.
    """
    img = Image.new("RGBA", (1, 1), (0, 0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.write(buf.getvalue())
    tmp.close()
    return Path(tmp.name)


def embed_audio(
    prs: Presentation,
    audio_map: dict[int, Path],
    slide_filter: set[int] | None,
    poster_frame: Path,
) -> int:
    """Embed WAV files into matching slides.

    Args:
        prs: Loaded Presentation object (modified in place).
        audio_map: Mapping of 1-based slide numbers to WAV file paths.
        slide_filter: Optional set of slide numbers to restrict embedding.
        poster_frame: Path to the poster frame image for add_movie.

    Returns:
        Count of slides that received embedded audio.
    """
    embedded_count = 0
    audio_top = prs.slide_height + AUDIO_OFFSCREEN_OFFSET
    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_filter and slide_num not in slide_filter:
            continue
        wav_path = audio_map.get(slide_num)
        if not wav_path:
            logger.debug("Slide %d: no audio file found, skipping", slide_num)
            continue

        # python-pptx does not expose a public audio-embedding API, so we use
        # add_movie which creates a video relationship type. PowerPoint Desktop
        # handles WAV media embedded this way correctly for narration timing and
        # video export via "Use Recorded Timings and Narrations". Other viewers
        # (LibreOffice, Google Slides) may display a video icon instead.
        slide.shapes.add_movie(
            movie_file=str(wav_path),
            left=AUDIO_LEFT,
            top=audio_top,
            width=AUDIO_WIDTH,
            height=AUDIO_HEIGHT,
            poster_frame_image=str(poster_frame),
            mime_type="audio/wav",
        )
        embedded_count += 1
        logger.info("Slide %d: embedded %s", slide_num, wav_path.name)

    return embedded_count


def run(args: argparse.Namespace) -> int:
    """Execute the audio embedding workflow.

    Args:
        args: Parsed command-line arguments.

    Returns:
        Exit code indicating success or failure.
    """
    input_path: Path = args.input
    audio_dir: Path = args.audio_dir
    output_path: Path = args.output

    if not input_path.is_file():
        logger.error("Input file not found: %s", input_path)
        return EXIT_ERROR

    if not audio_dir.is_dir():
        logger.error("Audio directory not found: %s", audio_dir)
        return EXIT_ERROR

    slide_filter = parse_slide_filter(args.slides)

    audio_map = discover_audio_files(audio_dir)
    if not audio_map:
        logger.warning("No slide-NNN.wav files found in %s", audio_dir)
        return EXIT_FAILURE

    logger.info("Discovered %d audio file(s) in %s", len(audio_map), audio_dir)

    prs = Presentation(str(input_path))
    total_slides = len(prs.slides)
    logger.info("Opened %s (%d slides)", input_path.name, total_slides)

    poster_frame = create_poster_frame()
    try:
        embedded = embed_audio(prs, audio_map, slide_filter, poster_frame)
    finally:
        poster_frame.unlink(missing_ok=True)

    if embedded == 0:
        logger.warning("No audio files matched any target slides")
        return EXIT_FAILURE

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("Saved %s with %d embedded audio track(s)", output_path, embedded)
    return EXIT_SUCCESS


def main() -> int:
    """Main entry point for the script."""
    parser = create_parser()
    args = parser.parse_args()
    configure_logging(args.verbose)
    try:
        return run(args)
    except KeyboardInterrupt:
        print("\nInterrupted by user", file=sys.stderr)
        return 130
    except BrokenPipeError:
        sys.stderr.close()
        return EXIT_FAILURE
    except Exception as e:
        logger.error("Unexpected error: %s", e)
        return EXIT_FAILURE


if __name__ == "__main__":
    sys.exit(main())
