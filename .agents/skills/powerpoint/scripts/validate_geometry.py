#!/usr/bin/env python3
# Copyright (c) 2026 Microsoft Corporation. All rights reserved.
# SPDX-License-Identifier: MIT
"""Validate PPTX element geometry against spacing and margin rules.

Checks edge margins, adjacent element gaps, boundary overflow, and
title-subtitle clearance. Decorative accent bars (full-width shapes at
top with height ≤ 0.12") are exempted from margin rules.

Usage::

    python validate_geometry.py --input deck.pptx
    python validate_geometry.py --input deck.pptx \
        --output results.json --report report.md
    python validate_geometry.py --input deck.pptx \
        --slides "1,3" --margin 0.6 --gap 0.4
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx_utils import (
    EXIT_ERROR,
    EXIT_FAILURE,
    EXIT_SUCCESS,
    configure_logging,
    emu_to_inches,
    parse_slide_filter,
)

logger = logging.getLogger(__name__)

SEVERITY_ICON = {"error": "❌", "warning": "⚠️", "info": "ℹ️"}
QUALITY_ICON = {"good": "✅", "needs-attention": "⚠️"}

ACCENT_BAR_MAX_HEIGHT = 0.12
POSITION_TOLERANCE_IN = 0.01  # floating-point tolerance for inch comparisons


def _is_accent_bar(shape: BaseShape, slide_width_in: float) -> bool:
    """Return True when shape is a full-width decorative accent bar at top."""
    top_in = emu_to_inches(shape.top)
    height_in = emu_to_inches(shape.height)
    width_in = emu_to_inches(shape.width)
    left_in = emu_to_inches(shape.left)
    return (
        top_in <= POSITION_TOLERANCE_IN
        and left_in <= POSITION_TOLERANCE_IN
        and height_in <= ACCENT_BAR_MAX_HEIGHT
        and abs(width_in - slide_width_in) < POSITION_TOLERANCE_IN
    )


def _is_offscreen_media(shape: BaseShape, slide_h_in: float) -> bool:
    """Return True when shape is placed entirely below the slide boundary.

    Audio shapes embedded by embed_audio.py are positioned off-screen
    below the slide height. These should not trigger boundary overflow errors.
    """
    top_in = emu_to_inches(shape.top)
    return top_in > slide_h_in


def _shape_label(shape: BaseShape) -> str:
    """Return a human-readable label for a shape."""
    name = shape.name or "unnamed"
    if hasattr(shape, "text") and shape.text:
        preview = shape.text[:40].replace("\n", " ")
        return f'{name} ("{preview}")'
    return name


def check_boundary_overflow(
    shape: BaseShape,
    slide_w_in: float,
    slide_h_in: float,
) -> list[dict]:
    """Check whether a shape extends beyond slide boundaries."""
    issues: list[dict] = []
    left = emu_to_inches(shape.left)
    top = emu_to_inches(shape.top)
    width = emu_to_inches(shape.width)
    height = emu_to_inches(shape.height)
    right = left + width
    bottom = top + height
    label = _shape_label(shape)
    if left < -POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "boundary_overflow",
                "severity": "error",
                "description": (
                    f"Shape '{label}' left edge ({left:.2f}\") extends "
                    "off the left boundary of the slide"
                ),
                "location": shape.name or "shape",
            }
        )
    if top < -POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "boundary_overflow",
                "severity": "error",
                "description": (
                    f"Shape '{label}' top edge ({top:.2f}\") extends "
                    "off the top boundary of the slide"
                ),
                "location": shape.name or "shape",
            }
        )
    if right > slide_w_in + POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "boundary_overflow",
                "severity": "error",
                "description": (
                    f"Shape '{label}' right edge ({right:.2f}\") exceeds "
                    f'slide width ({slide_w_in:.2f}")'
                ),
                "location": shape.name or "shape",
            }
        )
    if bottom > slide_h_in + POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "boundary_overflow",
                "severity": "error",
                "description": (
                    f"Shape '{label}' bottom edge ({bottom:.2f}\") exceeds "
                    f'slide height ({slide_h_in:.2f}")'
                ),
                "location": shape.name or "shape",
            }
        )
    return issues


def check_edge_margins(
    shape: BaseShape,
    slide_w_in: float,
    slide_h_in: float,
    margin: float,
) -> list[dict]:
    """Check whether a shape maintains minimum edge margins."""
    issues: list[dict] = []
    left = emu_to_inches(shape.left)
    top = emu_to_inches(shape.top)
    width = emu_to_inches(shape.width)
    height = emu_to_inches(shape.height)
    right = left + width
    bottom = top + height
    label = _shape_label(shape)

    if left < margin - POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "edge_margin",
                "severity": "warning",
                "description": (
                    f"Shape '{label}' left ({left:.2f}\") < minimum margin ({margin}\")"
                ),
                "location": shape.name or "shape",
            }
        )
    if top < margin - POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "edge_margin",
                "severity": "warning",
                "description": (
                    f"Shape '{label}' top ({top:.2f}\") < minimum margin ({margin}\")"
                ),
                "location": shape.name or "shape",
            }
        )
    if right > slide_w_in - margin + POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "edge_margin",
                "severity": "warning",
                "description": (
                    f"Shape '{label}' right edge ({right:.2f}\") > "
                    f'slide width - margin ({slide_w_in - margin:.2f}")'
                ),
                "location": shape.name or "shape",
            }
        )
    if bottom > slide_h_in - margin + POSITION_TOLERANCE_IN:
        issues.append(
            {
                "check_type": "edge_margin",
                "severity": "warning",
                "description": (
                    f"Shape '{label}' bottom edge ({bottom:.2f}\") > "
                    f'slide height - margin ({slide_h_in - margin:.2f}")'
                ),
                "location": shape.name or "shape",
            }
        )
    return issues


def check_adjacent_gaps(shapes: list[BaseShape], gap: float) -> list[dict]:
    """Check vertical gaps between vertically adjacent elements.

    Sorts shapes by top position and checks consecutive pairs for minimum
    vertical clearance. Only pairs that share horizontal extent are evaluated.

    Note: Horizontal gaps between side-by-side elements at the same vertical
    level are not checked by this function.
    """
    issues: list[dict] = []
    rects = []
    for s in shapes:
        top = emu_to_inches(s.top)
        height = emu_to_inches(s.height)
        rects.append((top, top + height, s))
    rects.sort(key=lambda r: r[0])

    for i in range(len(rects) - 1):
        _, bottom_a, shape_a = rects[i]
        top_b, _, shape_b = rects[i + 1]
        vertical_gap = top_b - bottom_a
        # Verify shapes share horizontal extent before flagging
        left_a = emu_to_inches(shape_a.left)
        right_a = left_a + emu_to_inches(shape_a.width)
        left_b = emu_to_inches(shape_b.left)
        right_b = left_b + emu_to_inches(shape_b.width)
        h_overlap = min(right_a, right_b) - max(left_a, left_b)
        too_close = vertical_gap < gap - POSITION_TOLERANCE_IN
        if too_close and vertical_gap >= 0 and h_overlap > POSITION_TOLERANCE_IN:
            label_a = _shape_label(shape_a)
            label_b = _shape_label(shape_b)
            issues.append(
                {
                    "check_type": "adjacent_gap",
                    "severity": "warning",
                    "description": (
                        f"Gap between '{label_a}' and '{label_b}' "
                        f'is {vertical_gap:.2f}" (minimum {gap}")'
                    ),
                    "location": f"{shape_a.name or 'shape'}→{shape_b.name or 'shape'}",
                }
            )
    return issues


def _is_title_placeholder(shape: BaseShape) -> bool:
    """Return True when shape is a title (not subtitle) placeholder.

    Prefers the placeholder type when available (robust across templates),
    falls back to shape name heuristic for non-placeholder shapes.
    """
    try:
        ph = shape.placeholder_format
        if ph is not None:
            # PP_PLACEHOLDER.TITLE = 1, CENTER_TITLE = 3; idx 0 is the
            # standard title placeholder index in most templates.
            return ph.idx == 0 or (ph.type is not None and ph.type in (1, 3))
    except (AttributeError, ValueError):
        # Non-placeholder shapes raise ValueError; fall through to name heuristic.
        pass
    # Fallback: name-based detection for non-placeholder shapes
    name_lower = (shape.name or "").lower()
    return "title" in name_lower and "subtitle" not in name_lower


def check_title_clearance(shapes: list[BaseShape], clearance: float) -> list[dict]:
    """Check title-to-next-element vertical clearance.

    Identifies title placeholders and verifies the next element below
    has sufficient clearance.
    """
    issues: list[dict] = []
    rects = []
    for s in shapes:
        top = emu_to_inches(s.top)
        height = emu_to_inches(s.height)
        rects.append((top, top + height, s))
    rects.sort(key=lambda r: r[0])

    for i, (_, bottom, shape) in enumerate(rects):
        if not _is_title_placeholder(shape):
            continue
        if i + 1 >= len(rects):
            continue
        next_top = rects[i + 1][0]
        title_clearance = next_top - bottom
        if title_clearance < clearance - POSITION_TOLERANCE_IN and title_clearance >= 0:
            label = _shape_label(shape)
            next_label = _shape_label(rects[i + 1][2])
            issues.append(
                {
                    "check_type": "title_clearance",
                    "severity": "info",
                    "description": (
                        f"Title '{label}' to '{next_label}' clearance "
                        f'is {title_clearance:.2f}" (recommended {clearance}")'
                    ),
                    "location": (
                        f"{shape.name or 'title'}→{rects[i + 1][2].name or 'shape'}"
                    ),
                }
            )
    return issues


def validate_slide_geometry(
    slide,
    slide_num: int,
    slide_w_in: float,
    slide_h_in: float,
    *,
    margin: float,
    gap: float,
    clearance: float,
) -> dict:
    """Run all geometry checks for a single slide."""
    issues: list[dict] = []
    non_accent_shapes = []

    for shape in slide.shapes:
        # Skip off-screen media shapes (e.g. audio embedded below slide boundary)
        if _is_offscreen_media(shape, slide_h_in):
            logger.debug(
                "Slide %d: exempting off-screen media '%s'",
                slide_num,
                shape.name,
            )
            continue

        # Boundary overflow applies to all visible shapes
        issues.extend(check_boundary_overflow(shape, slide_w_in, slide_h_in))

        if _is_accent_bar(shape, slide_w_in):
            logger.debug(
                "Slide %d: exempting accent bar '%s'",
                slide_num,
                shape.name,
            )
            continue

        non_accent_shapes.append(shape)
        issues.extend(check_edge_margins(shape, slide_w_in, slide_h_in, margin))

    # Adjacent gaps and title clearance use non-accent shapes only
    issues.extend(check_adjacent_gaps(non_accent_shapes, gap))
    issues.extend(check_title_clearance(non_accent_shapes, clearance))

    quality = "good" if not issues else "needs-attention"
    return {
        "slide_number": slide_num,
        "issues": issues,
        "overall_quality": quality,
    }


def validate_geometry(
    pptx_path: Path,
    slide_filter: set[int] | None = None,
    *,
    margin: float = 0.5,
    gap: float = 0.3,
    clearance: float = 0.2,
) -> dict:
    """Run geometry validation across all slides in a presentation.

    Returns:
        Dict with source, slide_count, and per-slide issues.
    """
    prs = Presentation(str(pptx_path))
    slide_w_in = emu_to_inches(prs.slide_width)
    slide_h_in = emu_to_inches(prs.slide_height)
    total_slides = len(prs.slides)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue
        slide_result = validate_slide_geometry(
            slide,
            slide_num,
            slide_w_in,
            slide_h_in,
            margin=margin,
            gap=gap,
            clearance=clearance,
        )
        slides.append(slide_result)

    return {
        "source": "geometry-validation",
        "slide_count": total_slides,
        "slides": slides,
    }


def generate_report(results: dict) -> str:
    """Generate a Markdown validation report from results."""
    lines = ["# Geometry Validation Report", ""]
    ts = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    lines.append(f"**Generated**: {ts}  ")
    lines.append(f"**Source**: {results['source']}  ")
    lines.append(f"**Slides**: {results['slide_count']}")
    lines.append("")

    error_count = 0
    warning_count = 0
    info_count = 0
    for slide in results["slides"]:
        for issue in slide.get("issues", []):
            sev = issue.get("severity", "info")
            if sev == "error":
                error_count += 1
            elif sev == "warning":
                warning_count += 1
            else:
                info_count += 1

    lines.append("## Summary")
    lines.append("")
    lines.append("| Severity | Count |")
    lines.append("|-|-|")
    lines.append(f"| ❌ Errors | {error_count} |")
    lines.append(f"| ⚠️ Warnings | {warning_count} |")
    lines.append(f"| ℹ️ Info | {info_count} |")
    lines.append("")

    lines.append("## Per-Slide Findings")
    lines.append("")
    for slide in results["slides"]:
        num = slide.get("slide_number", "?")
        quality = slide.get("overall_quality", "unknown")
        icon = QUALITY_ICON.get(quality, "❓")
        lines.append(f"### Slide {num} {icon} {quality}")
        lines.append("")

        issues = slide.get("issues", [])
        if not issues:
            lines.append("No issues found.")
            lines.append("")
            continue

        lines.append("| Severity | Check | Location | Description |")
        lines.append("|-|-|-|-|")
        for issue in issues:
            sev = issue.get("severity", "info")
            sev_icon = SEVERITY_ICON.get(sev, "")
            check = issue.get("check_type", "")
            loc = issue.get("location", "").replace("|", "\\|")
            desc = issue.get("description", "").replace("|", "\\|")
            lines.append(f"| {sev_icon} {sev} | {check} | {loc} | {desc} |")
        lines.append("")

    return "\n".join(lines)


def max_severity(results: dict) -> str:
    """Return the highest severity found across all issues."""
    severities = set()
    for slide in results["slides"]:
        for issue in slide.get("issues", []):
            severities.add(issue.get("severity", "info"))
    if "error" in severities:
        return "error"
    if "warning" in severities:
        return "warning"
    if "info" in severities:
        return "info"
    return "none"


def create_parser() -> argparse.ArgumentParser:
    """Create and configure argument parser."""
    parser = argparse.ArgumentParser(
        description=(
            "Validate PPTX element geometry: edge margins, adjacent gaps, "
            "boundary overflow, and title-subtitle clearance"
        )
    )
    parser.add_argument(
        "--input",
        required=True,
        type=Path,
        help="Input PPTX file path",
    )
    parser.add_argument(
        "--slides",
        help="Comma-separated slide numbers to validate (default: all)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        help="Output JSON file path (default: stdout)",
    )
    parser.add_argument(
        "--report",
        type=Path,
        help="Output Markdown report file path",
    )
    parser.add_argument(
        "--per-slide-dir",
        type=Path,
        help="Directory for per-slide JSON files (slide-NNN-geometry.json)",
    )
    parser.add_argument(
        "--margin",
        type=float,
        default=0.5,
        help="Minimum edge margin in inches (default: 0.5)",
    )
    parser.add_argument(
        "--gap",
        type=float,
        default=0.3,
        help="Minimum adjacent element gap in inches (default: 0.3)",
    )
    parser.add_argument(
        "--clearance",
        type=float,
        default=0.2,
        help="Minimum title-subtitle clearance in inches (default: 0.2)",
    )
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable verbose logging",
    )
    return parser


def run(args: argparse.Namespace) -> int:
    """Execute geometry validation logic."""
    pptx_path = args.input
    if not pptx_path.exists():
        logger.error("File not found: %s", pptx_path)
        return EXIT_ERROR

    if pptx_path.suffix.lower() != ".pptx":
        logger.error("Input file must be a .pptx file: %s", pptx_path)
        return EXIT_ERROR

    slide_filter = parse_slide_filter(args.slides)

    logger.info("Validating geometry: %s", pptx_path)
    results = validate_geometry(
        pptx_path,
        slide_filter=slide_filter,
        margin=args.margin,
        gap=args.gap,
        clearance=args.clearance,
    )

    # Write per-slide geometry JSON files
    if args.per_slide_dir:
        args.per_slide_dir.mkdir(parents=True, exist_ok=True)
        for slide_result in results["slides"]:
            slide_num = slide_result.get("slide_number", 0)
            per_slide_path = args.per_slide_dir / f"slide-{slide_num:03d}-geometry.json"
            per_slide_json = json.dumps(slide_result, indent=2)
            per_slide_path.write_text(per_slide_json, encoding="utf-8")
            logger.debug("Per-slide geometry results written to %s", per_slide_path)

    # Output JSON
    output_json = json.dumps(results, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(output_json, encoding="utf-8")
        logger.info("Results written to %s", args.output)
    else:
        print(output_json)

    # Generate Markdown report
    if args.report:
        report_md = generate_report(results)
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(report_md, encoding="utf-8")
        logger.info("Report written to %s", args.report)

    # Report summary
    total_issues = sum(len(s.get("issues", [])) for s in results["slides"])
    severity = max_severity(results)
    slide_count = results["slide_count"]
    logger.info(
        "Validation complete: %d issue(s) across %d slide(s)",
        total_issues,
        slide_count,
    )

    if severity == "error":
        return EXIT_ERROR
    if severity == "warning":
        return EXIT_FAILURE
    return EXIT_SUCCESS


def main() -> int:
    """Main entry point."""
    parser = create_parser()
    args = parser.parse_args()
    configure_logging(args.verbose)
    try:
        return run(args)
    except KeyboardInterrupt:
        return 130
    except BrokenPipeError:
        sys.stderr.close()
        return EXIT_FAILURE
    except Exception as e:
        logger.error("Unexpected error: %s", e)
        return EXIT_ERROR


if __name__ == "__main__":
    sys.exit(main())
